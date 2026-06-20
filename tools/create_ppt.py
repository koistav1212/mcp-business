import os
from typing import Optional, Type, List
from pydantic import BaseModel
from tools.base import BaseTool
from storage.local import local_storage
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

class SlideContent(BaseModel):
    title: str
    points: List[str]

class CreatePPTInput(BaseModel):
    filename: str
    presentation_title: str
    slides: List[SlideContent]
    style: Optional[str] = "professional"

class CreatePPTTool(BaseTool):
    name: str = "create_ppt"
    description: str = "Generates a structured PowerPoint presentation with slide lists."
    args_schema: Optional[Type[BaseModel]] = CreatePPTInput

    async def execute(self, **kwargs) -> str:
        filename = kwargs["filename"]
        if not filename.endswith(".pptx"):
            filename += ".pptx"
            
        presentation_title = kwargs["presentation_title"]
        slides_data = kwargs["slides"]
        style = kwargs.get("style", "professional")

        # RGB Colors
        if style == "creative":
            primary_color = RGBColor(99, 102, 241)  # Indigo
            text_color = RGBColor(79, 70, 229)     # Violet
        elif style == "minimalist":
            primary_color = RGBColor(9, 9, 11)      # Black
            text_color = RGBColor(39, 39, 42)      # Charcoal
        else: # professional
            primary_color = RGBColor(30, 41, 59)    # Dark Slate
            text_color = RGBColor(51, 65, 85)      # Slate Grey

        prs = Presentation()
        
        # 1. Add Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        # Color Title
        title_shape.text = presentation_title
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = primary_color

        subtitle_shape.text = "Business Intelligence Brief | Powered by AI Agent Framework"
        for paragraph in subtitle_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = text_color

        # 2. Add Bullet Content Slides
        bullet_layout = prs.slide_layouts[1]
        for slide_data in slides_data:
            # Handle list of dicts vs list of Pydantic models when executing
            if isinstance(slide_data, dict):
                slide_title = slide_data.get("title", "Info")
                slide_points = slide_data.get("points", [])
            else:
                slide_title = slide_data.title
                slide_points = slide_data.points

            slide = prs.slides.add_slide(bullet_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = slide_title
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = primary_color
            
            tf = body_shape.text_frame
            if slide_points:
                tf.text = slide_points[0]
                for pt in slide_points[1:]:
                    p = tf.add_paragraph()
                    p.text = pt
                    p.level = 0
                
                # Apply text color to body bullets
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = text_color
            else:
                tf.text = "No content details provided."
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = text_color

        # Save presentation
        temp_path = f"temp_gen_{filename}"
        prs.save(temp_path)

        with open(temp_path, "rb") as f:
            content = f.read()
            
        if os.path.exists(temp_path):
            os.remove(temp_path)

        download_url = local_storage.write_content(filename, content, "ppt")
        return download_url
